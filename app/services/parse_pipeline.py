"""Phase 2 — Parse pipeline orchestrating native + RAGFlow paths."""

from __future__ import annotations

import hashlib
import logging
import os
import re
import time
import asyncio
import csv
from typing import Any, Dict, List, Optional, Tuple
from uuid import UUID
from pathlib import Path

from app.schemas.parse_artifact import ParseArtifact, ParseChunk
from app.services.parse_router import ParseRoute, classify_document
from app.services.document_parser import DocumentParser
from app.services.content_reference import (
    build_content_reference,
    resolve_content_bytes,
)

logger = logging.getLogger(__name__)


def _update_heading_path(path: list[str], level: int, title: str) -> list[str]:
    """Return the heading ancestry after one real heading, never a list row."""
    normalized = str(title or "").strip()
    if not normalized or level < 1 or level > 6:
        return list(path)
    updated = list(path[: level - 1])
    while len(updated) < level - 1:
        updated.append("")
    updated.append(normalized)
    return [item for item in updated if item]


def _markdown_evidence_chunks(text: str) -> list[ParseChunk]:
    path: list[str] = []
    chunks: list[ParseChunk] = []
    heading_pattern = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
    for line in str(text or "").splitlines():
        value = line.strip()
        if not value:
            continue
        match = heading_pattern.match(value)
        if match:
            path = _update_heading_path(path, len(match.group(1)), match.group(2))
        chunks.append(
            ParseChunk(
                text=value,
                hierarchy=list(path),
                section=path[-1] if path else "document",
                chunk_index=len(chunks),
            )
        )
    return chunks


def _native_evidence_chunks(
    file_path: str,
    file_type: str,
    text: str,
    metadata: Dict[str, Any],
) -> List[ParseChunk]:
    """Preserve stable source coordinates for native parser output."""

    kind = (file_type or "").lower().lstrip(".")
    path = Path(file_path)
    try:
        if kind == "pdf":
            from pypdf import PdfReader

            pages = [
                str(page.extract_text() or "").strip()
                for page in PdfReader(file_path).pages
            ]
            chunks = [
                ParseChunk(text=value, page=index, chunk_index=index - 1)
                for index, value in enumerate(pages, start=1)
                if value
            ]
            if chunks:
                return chunks
            if text.strip() and len(pages) == 1:
                return [ParseChunk(text=text.strip(), page=1, chunk_index=0)]
        elif kind == "docx":
            from docx import Document as WordDocument

            document = WordDocument(file_path)
            sections: List[ParseChunk] = []
            heading_path: list[str] = []
            for paragraph_index, paragraph in enumerate(document.paragraphs, 1):
                value = paragraph.text.strip()
                if not value:
                    continue
                style_name = str(paragraph.style.name or "").strip().casefold()
                match = re.match(r"heading\s*([1-6])$", style_name)
                if match:
                    heading_path = _update_heading_path(
                        heading_path, int(match.group(1)), value
                    )
                sections.append(
                    ParseChunk(
                        text=value,
                        hierarchy=list(heading_path),
                        section=heading_path[-1] if heading_path else "document",
                        paragraph_index=paragraph_index,
                        chunk_index=len(sections),
                    )
                )
            for table_index, table in enumerate(document.tables, 1):
                for row_index, row in enumerate(table.rows, 1):
                    value = " | ".join(cell.text.strip() for cell in row.cells)
                    if value.strip(" |"):
                        sections.append(
                            ParseChunk(
                                text=value,
                                section=f"table:{table_index}:row:{row_index}",
                                chunk_index=len(sections),
                            )
                        )
            if sections:
                return sections
        elif kind in {"md", "markdown", "txt"}:
            chunks = _markdown_evidence_chunks(text)
            if chunks:
                return chunks
        elif kind in {"xlsx", "xls"}:
            from openpyxl import load_workbook
            from openpyxl.utils import get_column_letter

            workbook = load_workbook(file_path, read_only=False, data_only=False)
            chunks = []
            try:
                for sheet in workbook.worksheets:
                    if sheet.sheet_state != "visible":
                        continue
                    for row_number, row in enumerate(
                        sheet.iter_rows(values_only=False), sheet.min_row
                    ):
                        values = [
                            "" if cell.value is None else str(cell.value)
                            for cell in row
                        ]
                        if not any(value.strip() for value in values):
                            continue
                        start = get_column_letter(sheet.min_column)
                        end = get_column_letter(sheet.min_column + len(values) - 1)
                        chunks.append(
                            ParseChunk(
                                text=" | ".join(values),
                                worksheet=sheet.title,
                                row_number=row_number,
                                cell_range=f"{start}{row_number}:{end}{row_number}",
                                chunk_index=len(chunks),
                            )
                        )
            finally:
                workbook.close()
            if chunks:
                return chunks
        elif kind == "csv":
            with path.open("r", encoding="utf-8-sig", newline="") as stream:
                rows = list(csv.reader(stream))
            chunks = [
                ParseChunk(
                    text=" | ".join(row),
                    row_number=index,
                    chunk_index=index - 1,
                )
                for index, row in enumerate(rows, start=1)
                if any(str(value).strip() for value in row)
            ]
            if chunks:
                return chunks
        elif kind == "pptx":
            from pptx import Presentation

            chunks = []
            for slide_number, slide in enumerate(Presentation(file_path).slides, 1):
                values = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        values.extend(
                            paragraph.text.strip()
                            for paragraph in shape.text_frame.paragraphs
                            if paragraph.text.strip()
                        )
                    if getattr(shape, "has_table", False):
                        values.extend(
                            " | ".join(cell.text.strip() for cell in row.cells)
                            for row in shape.table.rows
                        )
                value = "\n".join(values).strip()
                if value:
                    chunks.append(
                        ParseChunk(
                            text=value,
                            page=slide_number,
                            slide_number=slide_number,
                            section=f"slide:{slide_number}",
                            chunk_index=len(chunks),
                        )
                    )
            if chunks:
                return chunks
        elif kind in {
            "jpg",
            "jpeg",
            "png",
            "tiff",
            "tif",
            "bmp",
            "webp",
            "heic",
            "image",
        }:
            declared = metadata.get("evidence_chunks") or []
            chunks = []
            for index, item in enumerate(declared):
                if (
                    not isinstance(item, dict)
                    or not str(item.get("text") or "").strip()
                ):
                    continue
                chunks.append(
                    ParseChunk(
                        text=str(item["text"]).strip(),
                        page=item.get("page"),
                        bbox=item.get("bbox"),
                        chunk_index=index,
                        locator_fallback=bool(item.get("locator_fallback", False)),
                    )
                )
            if chunks:
                return chunks
            return (
                [
                    ParseChunk(
                        text=text.strip(),
                        bbox={"x": 0.0, "y": 0.0, "w": 1.0, "h": 1.0},
                        chunk_index=0,
                        locator_fallback=True,
                    )
                ]
                if text.strip()
                else []
            )
    except Exception as exc:
        logger.warning("native evidence coordinate extraction failed: %s", exc)
    return (
        [ParseChunk(text=text.strip(), section="document", chunk_index=0)]
        if text.strip()
        else []
    )


class ScanParseDeliveryError(RuntimeError):
    """Scan/OCR route failed to deliver usable OCR text — must not complete ingest."""


def _content_hash(file_path: str) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return f"sha256:{h.hexdigest()}"


# RAGFlow's layout_recognize values map onto the engine label we are allowed to claim.
# "Plain Text" bypasses OCR/layout/TSR entirely, so it must never be labelled deepdoc.
_LAYOUT_ENGINE_MAP = {
    "deepdoc": ("ragflow/deepdoc", True),
    "plain text": ("ragflow/plaintext", False),
    "plaintext": ("ragflow/plaintext", False),
}


def _engine_label_for_layout(layout_recognize: Any) -> Tuple[str, bool]:
    """Return (engine_label, ocr_capable) for an upstream layout_recognize value."""
    if not layout_recognize:
        return "ragflow/unknown", False
    key = str(layout_recognize).strip().lower()
    if key in _LAYOUT_ENGINE_MAP:
        return _LAYOUT_ENGINE_MAP[key]
    # Vision-model layouts are named after the model; treat them as OCR-capable.
    return f"ragflow/{key.replace(' ', '_')}", True


def _looks_dirty_ocr(text: str) -> bool:
    """Heuristic for 'lots of chars but unusable OCR' (spaced CJK, garbage)."""
    t = (text or "").strip()
    if len(t) < 40:
        return True
    spaced_cjk = len(re.findall(r"[\u4e00-\u9fff]\s+[\u4e00-\u9fff]", t))
    if spaced_cjk >= 8:
        return True
    # High density of private-use / replacement-like noise
    noise = sum(1 for ch in t if ord(ch) in (0xFFFD,) or (0xE000 <= ord(ch) <= 0xF8FF))
    if len(t) > 0 and noise / len(t) > 0.05:
        return True
    # Low density of real letters/scripts vs punctuation soup (e.g. broken
    # Myanmar/Latin OCR that keeps length but loses readable text).
    if len(t) >= 80:
        useful = sum(
            1
            for ch in t
            if (
                "\u4e00" <= ch <= "\u9fff"  # CJK
                or "A" <= ch <= "Z"
                or "a" <= ch <= "z"
                or "\u1000" <= ch <= "\u109f"  # Myanmar
                or "\u0e00" <= ch <= "\u0e7f"  # Thai
                or "\u0400" <= ch <= "\u04ff"  # Cyrillic
                or "\u3040" <= ch <= "\u30ff"  # JP kana
                or "\uac00" <= ch <= "\ud7af"  # Hangul
            )
        )
        if useful / len(t) < 0.28:
            return True
        words = re.findall(r"[A-Za-z\u4e00-\u9fff\u1000-\u109f]{2,}", t)
        singles = re.findall(
            r"(?<![A-Za-z\u4e00-\u9fff])[A-Za-z](?![A-Za-z\u4e00-\u9fff])", t
        )
        if len(singles) >= 12 and (
            not words or len(singles) / max(len(words), 1) >= 1.5
        ):
            return True
    return False


def _scan_route(route: ParseRoute) -> bool:
    return route in (ParseRoute.RAGFLOW_DEEPDOC, ParseRoute.RAGFLOW_VLM)


def _cloud_ocr_should_run(
    file_type: str,
    text: str,
    artifact: ParseArtifact,
    route: ParseRoute,
) -> Tuple[bool, str]:
    """Return (should_run, reason). Expanded beyond 'too few chars'."""
    from app.services import cloud_ocr

    ext = (file_type or "").lower().strip()
    if not cloud_ocr.is_enabled() or ext not in cloud_ocr.SUPPORTED_EXTS:
        return False, "disabled_or_unsupported_ext"

    trigger = int(os.getenv("CLOUD_OCR_TRIGGER_MIN_CHARS", "200"))
    stripped = (text or "").strip()

    if artifact.parser == "native/text_fallback":
        return True, "text_fallback"
    if _scan_route(route) and not artifact.ocr_used:
        return True, "scan_route_without_ocr"
    if len(stripped) < trigger:
        return True, f"primary_yield_below_{trigger}_chars"
    if _looks_dirty_ocr(stripped):
        return True, "dirty_ocr_heuristic"
    return False, "not_needed"


def _enforce_scan_delivery(
    route: ParseRoute,
    artifact: ParseArtifact,
    text: str,
) -> None:
    """Refuse silent completion when DeepDOC/scan OCR did not deliver."""
    if os.getenv("SCAN_PARSE_STRICT", "true").lower() in ("0", "false", "no"):
        return
    if artifact.parser == "native/text_fallback":
        raise ScanParseDeliveryError(
            "scan parse undelivered: native/text_fallback "
            "(RAGFlow returned no usable chunks; cloud OCR did not rescue). "
            "Set SCAN_PARSE_STRICT=false to allow completed ingest (not recommended)."
        )
    if _scan_route(route) and not artifact.ocr_used:
        stripped = (text or "").strip()
        if len(stripped) < int(
            os.getenv("CLOUD_OCR_TRIGGER_MIN_CHARS", "200")
        ) or _looks_dirty_ocr(stripped):
            raise ScanParseDeliveryError(
                f"scan parse undelivered: route={route.value} parser={artifact.parser} "
                f"ocr_used=false with empty/dirty text"
            )


async def _parse_via_ragflow(
    file_path: str,
    file_type: str,
    document_id: UUID,
    revision: int,
    content_hash: str,
    route: ParseRoute,
    tenant_id: Optional[UUID] = None,
) -> ParseArtifact:
    from app.gateway.adapters.ragflow_http import RAGFlowHTTPAdapter
    from app.core.authorization import AuthorizationContext

    if tenant_id is not None:
        # ADR-013：生產路徑的 sidecar 歸屬以 tenant_sidecar_binding 為唯一權威
        from app.db.session import SessionLocal
        from app.services.sidecar_binding import resolve_ragflow_dataset_id

        db = SessionLocal()
        try:
            from app.services.rls import apply_rls_context

            apply_rls_context(db, tenant_id)
            dataset_id = resolve_ragflow_dataset_id(db, tenant_id)
        finally:
            db.close()
        if not dataset_id:
            raise RuntimeError(
                f"tenant {tenant_id} has no RAGFlow dataset bound — "
                "provision the pack before ingesting (ADR-013)"
            )
    else:
        # 無租戶上下文（測試／維運腳本）：部署級預設
        from app.services.sidecar_binding import legacy_env_dataset_id

        dataset_id = legacy_env_dataset_id() or ""
        if not dataset_id:
            raise RuntimeError(
                "RAGFLOW_DATASET_ID is empty — cannot run DeepDOC ingest. "
                "Configure the formal RAGFlow dataset id before scanning PDFs."
            )

    adapter = RAGFlowHTTPAdapter(
        base_url=os.getenv("RAGFLOW_BASE_URL", "http://localhost:9380"),
        api_key=os.getenv("RAGFLOW_API_KEY", ""),
    )
    authz = AuthorizationContext(
        tenant_id=UUID(int=0),
        subject_id=UUID(int=1),
        is_superuser=True,
    )
    start = time.time()
    content_ref, ref_meta = build_content_reference(file_path, UUID(int=0), document_id)
    file_bytes = resolve_content_bytes(content_ref, ref_meta)
    result = await adapter.ingest(
        document_id=document_id,
        revision=revision,
        content_uri=content_ref,
        content_hash=content_hash,
        file_type=file_type,
        authz=authz,
        metadata={
            "dataset_id": dataset_id,
            "ocr": route == ParseRoute.RAGFLOW_DEEPDOC,
            "vlm": route == ParseRoute.RAGFLOW_VLM,
            "file_path": file_path,
            "file_bytes": file_bytes,
        },
    )
    elapsed_ms = int((time.time() - start) * 1000)
    if result.get("status") == "error":
        raise RuntimeError(result.get("error", "ragflow ingest failed"))

    ragflow_doc_ids = result.get("ragflow_doc_ids") or []
    job_id = result.get("job_id") or (
        ragflow_doc_ids[0] if ragflow_doc_ids else "pending"
    )

    # DeepDOC on multi-page scans often exceeds 60s; default ~3 minutes.
    poll_attempts = int(os.getenv("RAGFLOW_PARSE_POLL_ATTEMPTS", "36"))
    poll_sleep = float(os.getenv("RAGFLOW_PARSE_POLL_SLEEP_S", "5"))
    parse_result: Dict[str, Any] = {"chunks": [], "confidence": 0.5}
    for _ in range(poll_attempts):
        parse_result = await adapter.get_parse_result(job_id)
        if parse_result.get("status") == "completed" and parse_result.get("chunks"):
            break
        if parse_result.get("status") == "error":
            break
        await asyncio.sleep(poll_sleep)
    chunks = []
    for i, c in enumerate(parse_result.get("chunks", [])):
        bbox_raw = c.get("bbox")
        bbox = None
        if isinstance(bbox_raw, dict):
            from app.schemas.parse_artifact import BBox

            bbox = BBox(
                x=float(bbox_raw.get("x", 0)),
                y=float(bbox_raw.get("y", 0)),
                w=float(bbox_raw.get("w", bbox_raw.get("width", 0))),
                h=float(bbox_raw.get("h", bbox_raw.get("height", 0))),
            )
        chunks.append(
            ParseChunk(
                text=c.get("text", ""),
                template=c.get("template", "general"),
                chunk_index=i,
                page=c.get("page") or c.get("page_num"),
                bbox=bbox,
                hierarchy=c.get("hierarchy") or [],
            )
        )
    # Label by what RAGFlow actually ran. The dataset's layout_recognize is the only
    # source of truth; the Enclave route only expresses intent.
    dataset_config = await adapter.get_dataset_config()
    layout_actual = (
        dataset_config.get("layout_recognize")
        if dataset_config.get("status") == "ok"
        else None
    )
    engine_label, layout_ocr_capable = _engine_label_for_layout(layout_actual)

    warnings: List[Any] = list(parse_result.get("warnings") or [])
    if dataset_config.get("status") != "ok":
        warnings.append(
            {
                "code": "ragflow_dataset_config_unavailable",
                "error": str(dataset_config.get("error"))[:200],
            }
        )

    raw_confidence = parse_result.get("confidence")
    confidence_provider_supplied = isinstance(
        raw_confidence, (int, float)
    ) and not isinstance(raw_confidence, bool)
    if not chunks:
        # RAGFlow produced nothing usable; label honestly as native fallback.
        # Delivery gate (+ cloud OCR rescue) decides whether ingest may complete.
        try:
            text, _meta = DocumentParser.parse(file_path, file_type)
        except Exception as exc:
            # Native fallback also failed (e.g. missing poppler makes the quality
            # gate reject the PDF). Do not let the raw error escape here — that
            # would bypass the cloud OCR rescue in _finish. Empty text keeps the
            # text_fallback label so the rescue/delivery gate decides the outcome.
            text = ""
            warnings.append(
                {
                    "code": "native_fallback_parse_failed",
                    "error": str(exc)[:300],
                }
            )
        chunks = [ParseChunk(text=text[:8000], chunk_index=0)]
        engine_label = "native/text_fallback"
        ocr_used = False
        # The provider score described an empty RAGFlow result, not the native
        # fallback text that replaces it.
        confidence = 0.5
        confidence_provider_supplied = False
        warnings.append({"code": "ragflow_chunks_empty_used_text_fallback"})
    else:
        ocr_used = layout_ocr_capable
        confidence = float(raw_confidence) if confidence_provider_supplied else 0.9

    return ParseArtifact(
        parser=engine_label,
        version=adapter.version,
        source_hash=content_hash,
        document_id=str(document_id),
        document_revision=revision,
        chunks=chunks,
        confidence=confidence,
        confidence_provider_supplied=confidence_provider_supplied,
        confidence_calibration_version=(
            "provider-native-uncalibrated"
            if confidence_provider_supplied
            else "parse-quality-heuristic.v1"
        ),
        elapsed_ms=elapsed_ms,
        ocr_used=ocr_used,
        vlm_used=route == ParseRoute.RAGFLOW_VLM
        and engine_label.startswith("ragflow/"),
        warnings=warnings,
        provider="ragflow",
        provider_resource_ids=[str(x) for x in ragflow_doc_ids if x],
        metadata={
            "layout_recognize_actual": layout_actual,
            "chunk_method_actual": dataset_config.get("chunk_method"),
            "parse_route_intent": route.value,
        },
    )


def _maybe_enhance_with_cloud_ocr(
    file_path: str,
    file_type: str,
    text: str,
    metadata: Dict[str, Any],
    artifact: ParseArtifact,
    route: ParseRoute,
) -> Tuple[str, Dict[str, Any], ParseArtifact]:
    """Cloud OCR enhancement arm (CV-RF-01b).

    Triggers on low yield, text_fallback, scan route without OCR, or dirty-OCR
    heuristic — not only 'too few characters'.
    """
    from app.services import cloud_ocr

    should, reason = _cloud_ocr_should_run(file_type, text, artifact, route)
    if not should:
        return text, metadata, artifact

    original_engine = artifact.parser
    try:
        results, candidate_errors = cloud_ocr.transcribe_candidates(
            file_path, (file_type or "").lower().strip()
        )
    except Exception as exc:
        artifact.warnings.append(
            {"code": "cloud_ocr_failed", "error": str(exc)[:200], "trigger": reason}
        )
        return text, metadata, artifact

    if not results:
        artifact.warnings.append(
            {
                "code": "cloud_ocr_failed",
                "error": "all configured OCR providers failed",
                "trigger": reason,
                "provider_errors": candidate_errors,
            }
        )
        return text, metadata, artifact

    def candidate_score(value: str) -> tuple[int, int, int]:
        stripped = (value or "").strip()
        useful = sum(character.isalnum() for character in stripped)
        return (0 if _looks_dirty_ocr(stripped) else 1, useful, len(stripped))

    result = max(results, key=lambda row: candidate_score(row.text))

    metadata["cloud_ocr"] = {
        "provider": result.provider,
        "model": result.model,
        "pages": result.pages,
        "elapsed_ms": result.elapsed_ms,
        "retries": result.retries,
        "errors": result.errors,
        "trigger": reason,
        "original_engine": original_engine,
        "candidate_count": len(results),
        "candidate_errors": candidate_errors,
        "candidates": [
            {
                "provider": row.provider,
                "model": row.model,
                "chars": len(row.text.strip()),
                "content_hash": hashlib.sha256(
                    row.text.encode("utf-8", errors="replace")
                ).hexdigest(),
                "selected": row is result,
            }
            for row in results
        ],
    }
    # Low-yield: require strictly more chars. Fallback/dirty/scan-without-ocr:
    # adopt non-empty cloud text that is cleaner or not shorter than half.
    cloud_len = len(result.text.strip())
    primary_len = len((text or "").strip())
    cloud_cleaner = (
        reason == "dirty_ocr_heuristic"
        and cloud_len >= 40
        and not _looks_dirty_ocr(result.text)
    )
    adopt = cloud_len > 0 and (
        cloud_len > primary_len
        or cloud_cleaner
        or (
            reason in ("text_fallback", "dirty_ocr_heuristic", "scan_route_without_ocr")
            and cloud_len >= max(40, primary_len // 2)
        )
    )
    if not adopt:
        artifact.warnings.append(
            {"code": "cloud_ocr_no_better_yield", "trigger": reason}
        )
        return text, metadata, artifact

    artifact.parser = f"cloud/{result.provider}:{result.model}"
    artifact.ocr_used = True
    # The cloud OCR response has no calibrated confidence. Never retain the
    # primary parser's score after replacing its content with another result.
    artifact.confidence = None
    artifact.confidence_provider_supplied = False
    artifact.confidence_calibration_version = "unavailable"
    artifact.chunks = [ParseChunk(text=result.text, chunk_index=0)]
    artifact.warnings.append(
        {
            "code": "cloud_ocr_adopted",
            "original_engine": original_engine,
            "original_chars": primary_len,
            "cloud_chars": cloud_len,
            "trigger": reason,
        }
    )
    metadata["parse_engine"] = artifact.parser
    metadata["ocr_used"] = True
    metadata["quality_score"] = None
    metadata["review_required"] = True
    return result.text, metadata, artifact


def _finish(
    file_path: str,
    file_type: str,
    text: str,
    metadata: Dict[str, Any],
    artifact: ParseArtifact,
    route: ParseRoute,
) -> Tuple[str, Dict[str, Any], ParseArtifact]:
    text, metadata, artifact = _maybe_enhance_with_cloud_ocr(
        file_path,
        file_type,
        text,
        metadata,
        artifact,
        route,
    )
    metadata["parse_engine"] = artifact.parser
    metadata["ocr_used"] = artifact.ocr_used
    _enforce_scan_delivery(route, artifact, text)
    return text, metadata, artifact


def parse_document(
    file_path: str,
    file_type: str,
    document_id: UUID,
    revision: int = 1,
    tenant_id: Optional[UUID] = None,
) -> Tuple[str, Dict[str, Any], ParseArtifact]:
    """
    Parse document and return (text_content, metadata_dict, ParseArtifact).

    Scan/DeepDOC routes must deliver OCR-capable text. native/text_fallback is
    labelled honestly but rejected by default (SCAN_PARSE_STRICT=true) unless
    cloud OCR rescues the document.

    ``tenant_id``：生產路徑（document_tasks）必帶，RAGFlow 歸屬走
    tenant_sidecar_binding（ADR-013）；None 僅限測試／維運腳本。
    """
    content_hash = _content_hash(file_path)
    route = classify_document(file_path, file_type)
    start = time.time()
    force_ragflow = os.getenv("RAGFLOW_FORCE_PARSE", "").lower() == "true"
    ragflow_error: Optional[str] = None

    if (
        route in (ParseRoute.RAGFLOW_DEEPDOC, ParseRoute.RAGFLOW_VLM)
        and os.getenv("RAGFLOW_ENABLED", "").lower() == "true"
    ):
        try:
            artifact = asyncio.run(
                _parse_via_ragflow(
                    file_path,
                    file_type,
                    document_id,
                    revision,
                    content_hash,
                    route,
                    tenant_id=tenant_id,
                )
            )
            text = "\n\n".join(c.text for c in artifact.chunks if c.text)
            metadata = {
                "parse_engine": artifact.parser,
                "parse_route": route.value,
                "quality_score": artifact.confidence,
                "ocr_used": artifact.ocr_used,
                "vlm_used": artifact.vlm_used,
                "content_hash": content_hash,
                "elapsed_ms": artifact.elapsed_ms,
                "ragflow_already_ingested": bool(artifact.provider_resource_ids),
                "ragflow_doc_ids": list(artifact.provider_resource_ids),
                "layout_recognize_actual": artifact.metadata.get(
                    "layout_recognize_actual"
                ),
                "chunk_method_actual": artifact.metadata.get("chunk_method_actual"),
            }
            return _finish(file_path, file_type, text, metadata, artifact, route)
        except ScanParseDeliveryError:
            raise
        except Exception as exc:
            if force_ragflow:
                raise RuntimeError(f"RAGFlow parse required but failed: {exc}") from exc
            logger.warning(
                "RAGFlow parse failed, attempting native+cloud rescue: %s", exc
            )
            ragflow_error = str(exc)[:300]

    # Structured tables (xlsx/csv/xls) are intentionally routed to native even when
    # RAGFLOW_FORCE_PARSE / PARSER_CANARY=ragflow — do not treat that as a failure.
    if (
        force_ragflow
        and os.getenv("RAGFLOW_ENABLED", "").lower() == "true"
        and route not in (ParseRoute.NATIVE_STRUCTURED,)
        and ragflow_error is None
    ):
        raise RuntimeError(
            "RAGFLOW_FORCE_PARSE=true but document was not routed to RAGFlow"
        )

    warnings: List[Any] = []
    # P3-4：Docling Parser（feature-flagged，條件式採用）
    docling_text = ""
    docling_metadata: Dict[str, Any] = {}
    from app.config import settings

    if settings.DOCLING_ENABLED:
        from app.services.docling_ablation import DoclingParser

        docling = DoclingParser()
        if docling.is_available():
            try:
                docling_result = docling.parse(file_path, file_type)
                if docling_result.success and docling_result.text:
                    docling_text = docling_result.text
                    docling_metadata = {
                        "docling_used": True,
                        "docling_elapsed_ms": int(
                            docling_result.elapsed_seconds * 1000
                        ),
                        "docling_tables": len(docling_result.tables),
                    }
            except Exception as exc:
                logger.warning("Docling parse failed: %s", exc)

    try:
        text_content, metadata = DocumentParser.parse(file_path, file_type)
        # 若 Docling 產出更多文字，採用 Docling
        if docling_text and len(docling_text) > len(text_content or "") * 1.2:
            text_content = docling_text
            metadata.update(docling_metadata)
            metadata["parse_engine"] = "docling"
            metadata["docling_adopted"] = True
            metadata["quality_score"] = None
            metadata["review_required"] = True
    except ValueError as exc:
        if not _scan_route(route):
            raise
        # Scan route but native quality gate rejected the document (e.g. missing
        # poppler). Keep the honest text_fallback label and let _finish give the
        # cloud OCR rescue its chance; the delivery gate fails the doc with an
        # actionable ScanParseDeliveryError if rescue is unavailable.
        text_content = ""
        metadata = {
            "parse_engine": "native/text_fallback",
            "quality_score": 0.0,
            "native_error": str(exc)[:300],
        }
        warnings.append(
            {"code": "native_parse_quality_rejected", "error": str(exc)[:300]}
        )
    if ragflow_error is not None:
        metadata["ragflow_error"] = ragflow_error
        warnings.append({"code": "ragflow_exception_native_rescue"})
        if _scan_route(route):
            metadata["parse_engine"] = "native/text_fallback"
    elapsed_ms = int((time.time() - start) * 1000)
    engine = metadata.get("parse_engine", "native")
    quality_score = metadata.get("quality_score", 0.8)
    confidence = (
        float(quality_score)
        if isinstance(quality_score, (int, float))
        and not isinstance(quality_score, bool)
        else None
    )
    artifact = ParseArtifact(
        parser=engine,
        version="1.0.0",
        source_hash=content_hash,
        document_id=str(document_id),
        document_revision=revision,
        chunks=_native_evidence_chunks(file_path, file_type, text_content, metadata),
        confidence=confidence,
        confidence_provider_supplied=False,
        confidence_calibration_version=(
            "parse-quality-heuristic.v1" if confidence is not None else "unavailable"
        ),
        elapsed_ms=elapsed_ms,
        ocr_used=bool(metadata.get("ocr_used", False)),
        warnings=warnings,
    )
    metadata["parse_route"] = (
        route.value if _scan_route(route) else ParseRoute.NATIVE_FAST.value
    )
    metadata["content_hash"] = content_hash
    return _finish(file_path, file_type, text_content, metadata, artifact, route)
