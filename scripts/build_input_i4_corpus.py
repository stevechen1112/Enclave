#!/usr/bin/env python3
"""Build and seal the synthetic manufacturing regression corpus for Input I4.

The corpus contains no customer data.  It exercises the structures and image
conditions found in SOPs, BOMs, inspection decks, scans and nameplates.  Real
plant/customer validation remains a separate, explicitly declared gate.
"""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageEnhance, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "artifacts" / "input" / "i4_corpus"
MANIFEST = ROOT / "artifacts" / "input" / "i4_quality_corpus_manifest.json"


def _font(size: int):
    candidates = (
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        Path("C:/Windows/Fonts/arialbd.ttf"),
    )
    for candidate in candidates:
        if candidate.is_file():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def _label(lines: list[str], *, dark: bool = False) -> Image.Image:
    background = (38, 42, 46) if dark else (248, 248, 242)
    foreground = (115, 115, 110) if dark else (18, 18, 18)
    image = Image.new("RGB", (1400, 620), background)
    draw = ImageDraw.Draw(image)
    for index, line in enumerate(lines):
        draw.text((80, 70 + index * 125), line, font=_font(64), fill=foreground)
    return image


def _save_docx(path: Path) -> None:
    document = Document()
    document.add_heading("CNC 換線標準作業", level=1)
    document.add_paragraph("適用設備：CNC-M2048")
    document.add_heading("停機確認", level=2)
    document.add_paragraph("確認主軸完全停止，壓力表歸零後再開啟安全門。")
    table = document.add_table(rows=2, cols=3)
    for cell, value in zip(table.rows[0].cells, ("步驟", "責任角色", "判定")):
        cell.text = value
    for cell, value in zip(table.rows[1].cells, ("鎖定能源", "操作員", "LOTO-OK")):
        cell.text = value
    document.save(path)


def _save_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "BOM"
    sheet.append(["料號", "品名", "單機用量", "批量", "總需求"])
    sheet.append(["AX-2048", "主軸護蓋", 2, 6, "=C2*D2"])
    sheet.append(["BR-7710", "耐油軸承", 4, 6, "=C3*D3"])
    sheet.merge_cells("A5:E5")
    sheet["A5"] = "工程變更 ECN-104 生效後使用"
    hidden = workbook.create_sheet("成本_機密")
    hidden.sheet_state = "hidden"
    hidden.append(["料號", "成本"])
    hidden.append(["AX-2048", 999])
    workbook.save(path)


def _save_pptx(path: Path) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "首件檢驗作業"
    box = first.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    box.text_frame.text = "量測特性：外徑 OD-42.00 mm\n允收範圍：41.98–42.02 mm"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "異常隔離"
    box = second.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    box.text_frame.text = "超差時標記 NCR-STOP，禁止混入合格品。"
    presentation.save(path)


def _entry(path: Path, *, entry_id: str, category: str, condition: str, expectations: list[str], positive: bool = True) -> dict:
    payload = path.read_bytes()
    return {
        "id": entry_id,
        "path": path.relative_to(ROOT).as_posix(),
        "sha256": hashlib.sha256(payload).hexdigest(),
        "bytes": len(payload),
        "extension": path.suffix.lower(),
        "category": category,
        "condition": condition,
        "positive_control": positive,
        "expectations": expectations,
    }


def main() -> int:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    procedure_txt = OUTPUT / "maintenance_procedure.txt"
    inspection_csv = OUTPUT / "incoming_inspection.csv"
    docx = OUTPUT / "cnc_changeover_sop.docx"
    xlsx = OUTPUT / "assembly_bom.xlsx"
    pptx = OUTPUT / "first_article_inspection.pptx"
    nameplate = OUTPUT / "equipment_nameplate.png"
    nameplate_jpg = OUTPUT / "equipment_nameplate.jpg"
    nameplate_jpeg = OUTPUT / "equipment_nameplate.jpeg"
    rotated = OUTPUT / "rotated_part_label.png"
    lowlight = OUTPUT / "lowlight_inspection_tag.png"
    tiff = OUTPUT / "multipage_scan_form.tiff"
    procedure_txt.write_text(
        "設備保養程序 MNT-2048\n1. 關閉 CNC-M2048 主電源。\n2. 確認壓力歸零。\n3. 完成點檢後填寫 WO-829。\n",
        encoding="utf-8",
    )
    inspection_csv.write_text(
        "檢驗批號,料號,抽樣數,結果\nIQC-240829,AX-2048,8,ACCEPT\nIQC-240830,BR-7710,8,HOLD\n",
        encoding="utf-8",
    )
    _save_docx(docx)
    _save_xlsx(xlsx)
    _save_pptx(pptx)
    clean_label = _label(["EQUIPMENT CNC-M2048", "SERIAL SN-771042", "MAX PRESSURE 6.5 BAR"])
    clean_label.save(nameplate)
    clean_label.save(nameplate_jpg, quality=94, subsampling=0)
    clean_label.save(nameplate_jpeg, quality=94, subsampling=0)
    _label(["PART AX-2048", "LOT L240829", "QTY 120"]).rotate(90, expand=True).save(rotated)
    dark = _label(["INSPECTION IPQC-204", "RESULT HOLD", "NCR NCR-STOP"], dark=True)
    dark = ImageEnhance.Brightness(dark).enhance(0.62)
    dark.save(lowlight)
    page_one = _label(["FORM IQC-204", "PART BR-7710", "SAMPLE 8"])
    page_two = _label(["RESULT ACCEPT", "REVIEWER OP-07", "REV 3"])
    page_one.save(tiff, save_all=True, append_images=[page_two], compression="tiff_deflate")

    entries = [
        _entry(procedure_txt, entry_id="maintenance-txt", category="maintenance_procedure", condition="utf8_multilingual", expectations=["MNT-2048", "CNC-M2048", "WO-829"]),
        _entry(inspection_csv, entry_id="inspection-csv", category="incoming_inspection", condition="row_cell", expectations=["IQC-240829", "AX-2048", "HOLD"]),
        _entry(ROOT / "tests/fixtures/nas_share/quality_manual.pdf", entry_id="manual-pdf", category="equipment_manual", condition="text_page", expectations=["Manufacturing Quality Manual", "Incoming inspection"]),
        _entry(docx, entry_id="sop-docx", category="sop", condition="headings_table", expectations=["CNC-M2048", "LOTO-OK"]),
        _entry(xlsx, entry_id="bom-xlsx", category="bom", condition="formula_merged_hidden", expectations=["AX-2048", "=C2*D2", "ECN-104"]),
        _entry(pptx, entry_id="inspection-pptx", category="inspection", condition="multi_slide", expectations=["OD-42.00", "NCR-STOP"]),
        _entry(nameplate, entry_id="nameplate-png", category="nameplate", condition="clean", expectations=["CNC-M2048", "SN-771042", "6.5 BAR"]),
        _entry(nameplate_jpg, entry_id="nameplate-jpg", category="nameplate", condition="clean_jpeg", expectations=["CNC-M2048", "SN-771042", "6.5 BAR"]),
        _entry(nameplate_jpeg, entry_id="nameplate-jpeg", category="nameplate", condition="clean_jpeg_alias", expectations=["CNC-M2048", "SN-771042", "6.5 BAR"]),
        _entry(rotated, entry_id="rotated-label-png", category="part_label", condition="rotated_90", expectations=["AX-2048", "L240829"]),
        _entry(lowlight, entry_id="lowlight-tag-png", category="inspection_tag", condition="low_light", expectations=["IPQC-204", "NCR-STOP"]),
        _entry(tiff, entry_id="multipage-tiff", category="scanned_form", condition="two_page", expectations=["IQC-204", "BR-7710", "ACCEPT", "OP-07"]),
    ]
    manifest = {
        "schema_version": 1,
        "phase": "Input I4",
        "status": "SEALED_INTERNAL_SYNTHETIC",
        "frozen_at": "2026-08-29T00:00:00Z",
        "purpose": "Manufacturing-shaped parser quality regression; not customer or field certification",
        "entries": entries,
        "declared_gaps": [
            "real handwritten shop-floor annotations",
            "customer-owned legacy DOC/XLS samples",
            "physical-device HEIC variants",
            "real damaged and oily paper samples",
            "provider production drift replay requires configured credentials",
        ],
    }
    encoded_entries = json.dumps(entries, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")
    manifest["corpus_sha256"] = hashlib.sha256(encoded_entries).hexdigest()
    MANIFEST.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(MANIFEST)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
