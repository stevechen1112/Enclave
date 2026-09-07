from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.schemas.parse_artifact import ParseChunk
from app.services.document_parser import DocumentParser, QualityReport
from app.services.input_quality import (
    content_accuracy,
    evaluate_observations,
    provider_drift,
    requires_human_review,
)
from app.services.parse_pipeline import _native_evidence_chunks


def test_quality_evaluator_measures_content_and_locator_not_only_parse_success():
    result = evaluate_observations(
        ".xlsx",
        [
            {
                "id": "bom-part-number",
                "parse_success": True,
                "expected": "AX-2048",
                "actual": "AX-2O48",
                "locator_complete": True,
            }
        ],
    )
    assert result["parse_success"] == 1.0
    assert result["status"] == "FAIL"
    assert result["content_accuracy"] < 1.0
    assert result["failures"][0]["id"] == "bom-part-number"


def test_content_accuracy_accepts_expected_key_field_inside_source_context():
    assert content_accuracy("M-2048", "設備編號：M-2048；線別：A3") == 1.0


def test_provider_drift_fails_closed_on_material_regression():
    baseline = {
        "parse_success": 1.0,
        "content_accuracy": 1.0,
        "locator_coverage": 1.0,
    }
    candidate = {**baseline, "content_accuracy": 0.90}
    assert provider_drift(".docx", baseline=baseline, candidate=candidate)["status"] == "FAIL"


def test_review_sampling_is_deterministic_and_fallback_always_reviews():
    first = requires_human_review(".docx", confidence=0.99, content_hash="same")
    second = requires_human_review(".docx", confidence=0.99, content_hash="same")
    assert first is second
    assert requires_human_review(
        ".png", confidence=0.99, content_hash="high", fallback_used=True
    )


def test_docx_evidence_preserves_paragraph_index(tmp_path: Path):
    path = tmp_path / "sop.docx"
    document = Document()
    document.add_heading("換線作業", level=1)
    document.add_paragraph("確認主電源已關閉")
    document.save(path)

    text, metadata = DocumentParser.parse(str(path), "docx")
    chunks = _native_evidence_chunks(str(path), "docx", text, metadata)
    assert [chunk.paragraph_index for chunk in chunks] == [1, 2]
    assert chunks[1].section == "換線作業"


def test_xlsx_preserves_formula_row_cell_and_excludes_hidden_sheet(tmp_path: Path):
    path = tmp_path / "bom.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "BOM"
    sheet.append(["料號", "數量", "總量"])
    sheet.append(["AX-2048", 2, "=B2*3"])
    sheet.merge_cells("A4:C4")
    sheet["A4"] = "核准後生效"
    hidden = workbook.create_sheet("成本")
    hidden.sheet_state = "hidden"
    hidden.append(["機密", 999])
    workbook.save(path)

    text, metadata = DocumentParser.parse(str(path), "xlsx")
    chunks = _native_evidence_chunks(str(path), "xlsx", text, metadata)
    assert "=B2*3" in text
    assert metadata["structure_policy"]["hidden_sheets"] == ["成本"]
    assert metadata["structure_policy"]["merged_cells_detected"] == 1
    assert metadata["structure_policy"]["merged_ranges"] == {"BOM": ["A4:C4"]}
    assert metadata["structure_policy"]["formula_cells"] == [
        {"worksheet": "BOM", "cell": "C2", "formula": "=B2*3"}
    ]
    assert {chunk.worksheet for chunk in chunks} == {"BOM"}
    assert next(chunk for chunk in chunks if "AX-2048" in chunk.text).cell_range == "A2:C2"
    assert next(chunk for chunk in chunks if "AX-2048" in chunk.text).row_number == 2


def test_xlsx_nonstandard_drawing_xml_uses_read_only_cell_fallback(
    tmp_path: Path, monkeypatch
):
    path = tmp_path / "legacy-erp-export.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "檢驗申請"
    sheet.append(["委託單號", "品名", "數量"])
    sheet.append(["QA-2026-0907", "草本原料", "12"])
    workbook.save(path)

    import openpyxl

    real_load_workbook = openpyxl.load_workbook
    calls: list[bool] = []

    def compatibility_only(source, *, read_only=False, data_only=False, **kwargs):
        calls.append(read_only)
        if not read_only:
            raise ValueError("drawing XML pitchFamily is outside supported range")
        return real_load_workbook(
            source,
            read_only=True,
            data_only=data_only,
            **kwargs,
        )

    monkeypatch.setattr(openpyxl, "load_workbook", compatibility_only)

    text, metadata = DocumentParser.parse(str(path), "xlsx")
    chunks = _native_evidence_chunks(str(path), "xlsx", text, metadata)

    assert calls == [False, True, False, True]
    assert "QA-2026-0907" in text
    assert metadata["parse_engine"] == "native/openpyxl-read-only"
    assert metadata["structure_policy"]["compatibility_mode"] is True
    assert (
        metadata["structure_policy"]["merged_cell_policy"]
        == "unavailable_in_read_only_compatibility"
    )
    assert any("非標準繪圖或 XML" in warning for warning in metadata["warnings"])
    row = next(chunk for chunk in chunks if "QA-2026-0907" in chunk.text)
    assert row.worksheet == "檢驗申請"
    assert row.row_number == 2
    assert row.cell_range == "A2:C2"


def test_pptx_evidence_preserves_slide_number(tmp_path: Path):
    path = tmp_path / "inspection.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "首件檢驗"
    presentation.save(path)

    text, metadata = DocumentParser.parse(str(path), "pptx")
    chunks = _native_evidence_chunks(str(path), "pptx", text, metadata)
    assert chunks[0].page == 1
    assert chunks[0].slide_number == 1
    assert "首件檢驗" in chunks[0].text


def test_image_declared_regions_beat_whole_image_fallback():
    metadata = {
        "evidence_chunks": [
            {
                "text": "M-2048",
                "bbox": {"x": 0.1, "y": 0.2, "w": 0.3, "h": 0.1},
            }
        ]
    }
    chunks = _native_evidence_chunks("missing.png", "png", "M-2048", metadata)
    assert chunks == [
        ParseChunk(
            text="M-2048",
            bbox={"x": 0.1, "y": 0.2, "w": 0.3, "h": 0.1},
            chunk_index=0,
            locator_fallback=False,
        )
    ]


def test_quality_report_exposes_explicit_structure_and_fallback_fields():
    payload = QualityReport().to_dict()
    assert payload["structure_policy"] == {}
    assert payload["evidence_chunks"] == []
    assert payload["locator_fallback"] is False
